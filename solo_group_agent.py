#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Groupwork Ultra Agent (Extreme Edition)
=======================================

What it is
----------
A local "groupwork OS" that lets a single student simulate a full group process
using multiple LLMs (different models) + a controller model that decides when
the discussion is "good enough" to stop.

Key features
------------
- SQLite persistence + Flask web app (vote & signup), inspired by your reference program.
- Multi-model, multi-role deliberation loops (Leader, Critic, Researcher, Methodologist,
  Editor, Red-Teamer, etc.) with merge & quality gates.
- A separate CoT Controller model that decides whether more rounds are worth it.
  (It reasons internally but outputs only short JSON verdict.)
- Local document ingestion from:
    ./Background_Information
    ./Assignment_Requirement
  Supports many common formats:
    .txt .md .docx .pptx .pdf
    .doc .ppt via LibreOffice 'soffice' conversion (optional)
- Local retrieval (RAG) via SQLite FTS5 chunks, to reduce hallucinations and keep
  responses grounded in your provided materials.
- Autopilot mode (AI picks topic & tasks) and Human-in-the-loop mode (web voting/signup).
- Robust JSON-mode chat completion with auto-repair retries.

Ethics / policy note
--------------------
This tool is designed for planning, structuring, and summarizing based on provided materials.
You should verify facts, add your own analysis, and comply with your course policy on AI usage.
"""

from __future__ import annotations

import argparse
import datetime as dt
import hashlib
import json
import os
import random
import re
import sqlite3
import subprocess
import tempfile
import textwrap
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from dotenv import load_dotenv
load_dotenv(dotenv_path=".env", override=False)

import requests

try:
    from flask import Flask, abort, redirect, render_template_string, request, url_for
except Exception:  # pragma: no cover
    Flask = None  # type: ignore


# ----------------------------- Defaults ---------------------------------------

DEFAULT_DB = "groupwork_ultra.sqlite3"
DEFAULT_HOST = "127.0.0.1"
DEFAULT_PORT = 8080

DEFAULT_API_BASE = "https://openrouter.ai/api/v1"
DEFAULT_MODELS = {
    "leader": "openai/gpt-5.2",
    "critic": "x-ai/grok-4.1-fast",
    "researcher": "google/gemini-3-pro-preview",
    "methodologist": "openai/gpt-5.2",
    "editor": "google/gemini-3-pro-preview",
    "redteam": "x-ai/grok-4.1-fast",
    "controller": "openai/gpt-5.2",
}

DEFAULT_BG_DIR = "./Background_Information"
DEFAULT_REQ_DIR = "./Assignment_Requirement"


# ----------------------------- Utilities --------------------------------------

print("RUNNING FILE:", __file__)
def utc_now_iso() -> str:
    return dt.datetime.utcnow().replace(microsecond=0).isoformat() + "Z"


def sha256_hex(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8")).hexdigest()


def ensure_dir(path: str) -> None:
    os.makedirs(path, exist_ok=True)


def choose_api_key(explicit_key: Optional[str]) -> str:
    if explicit_key:
        return explicit_key
    for k in ("LLM_API_KEY", "OPENROUTER_API_KEY", "OPENAI_API_KEY"):
        v = os.getenv(k)
        if v:
            return v
    raise SystemExit(
        "No API key found. Pass --api-key or set env: LLM_API_KEY / OPENROUTER_API_KEY / OPENAI_API_KEY"
    )


def try_parse_json(raw: str) -> Dict[str, Any]:
    raw = raw.strip()
    try:
        return json.loads(raw)
    except json.JSONDecodeError:
        # best-effort extraction between first { and last }
        i = raw.find("{")
        j = raw.rfind("}")
        if i != -1 and j != -1 and j > i:
            snippet = raw[i : j + 1]
            return json.loads(snippet)
        raise


def clamp(n: int, lo: int, hi: int) -> int:
    return max(lo, min(hi, n))


# ----------------------------- Multi-format reading ---------------------------

def _read_text_file(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()


def _read_docx(path: str) -> str:
    try:
        import docx  # type: ignore
    except ImportError:
        return ""
    doc = docx.Document(path)
    out: List[str] = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if t:
            out.append(t)
    return "\n".join(out)


def _read_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ""
    prs = Presentation(path)
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    texts.append(t)
    return "\n\n".join(texts)


def _read_pdf(path: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return ""
    chunks: List[str] = []
    with open(path, "rb") as f:
        reader = PdfReader(f)
        for page in reader.pages:
            t = (page.extract_text() or "").strip()
            if t:
                chunks.append(t)
    return "\n\n".join(chunks)



def _soffice_convert_to_txt(path: str) -> str:
    """
    Convert legacy .doc/.ppt to txt using LibreOffice 'soffice'.
    Returns empty string if not available or conversion fails.
    """
    soffice = shutil_which("soffice")
    if not soffice:
        return ""
    p = Path(path)
    with tempfile.TemporaryDirectory() as tmpdir:
        cmd = [
            soffice,
            "--headless",
            "--convert-to",
            "txt:Text",
            "--outdir",
            tmpdir,
            path,
        ]
        proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=False)
        if proc.returncode != 0:
            return ""
        txt_path = os.path.join(tmpdir, p.stem + ".txt")
        if not os.path.exists(txt_path):
            return ""
        return _read_text_file(txt_path)


def shutil_which(cmd: str) -> Optional[str]:
    # local replacement to avoid importing shutil late; keep single-file minimal
    for p in os.getenv("PATH", "").split(os.pathsep):
        full = os.path.join(p, cmd)
        if os.name == "nt":
            if os.path.exists(full) or os.path.exists(full + ".exe"):
                return full if os.path.exists(full) else full + ".exe"
        else:
            if os.path.exists(full) and os.access(full, os.X_OK):
                return full
    return None


def extract_text_from_file(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext in {".txt", ".md"}:
        return _read_text_file(path)
    if ext == ".docx":
        return _read_docx(path)
    if ext == ".pptx":
        return _read_pptx(path)
    if ext == ".pdf":
        return _read_pdf(path)
    if ext in {".doc", ".ppt"}:
        return _soffice_convert_to_txt(path)
    return ""


def load_directory_documents(directory: str) -> List[Tuple[str, str]]:
    """
    Returns list of (filename, extracted_text) for supported files in directory.
    Skips unreadable or empty results.
    """
    if not os.path.isdir(directory):
        raise FileNotFoundError(f"Directory not found: {directory}")
    out: List[Tuple[str, str]] = []
    for name in sorted(os.listdir(directory)):
        if name.startswith("."):
            continue
        path = os.path.join(directory, name)
        if not os.path.isfile(path):
            continue
        text = extract_text_from_file(path).strip()
        if text:
            out.append((name, text))
    return out


# ----------------------------- DB (SQLite + FTS) ------------------------------

class DB:
    def __init__(self, path: str):
        self.path = path
        self.conn = sqlite3.connect(path, check_same_thread=False)
        self.conn.row_factory = sqlite3.Row
        self.init_schema()

    def init_schema(self) -> None:
        cur = self.conn.cursor()
        cur.executescript(
            """
            PRAGMA journal_mode=WAL;

            CREATE TABLE IF NOT EXISTS projects (
              id TEXT PRIMARY KEY,
              name TEXT NOT NULL,
              created_at TEXT NOT NULL,
              leader TEXT,
              status TEXT NOT NULL DEFAULT 'init',
              settings_json TEXT NOT NULL DEFAULT '{}'
            );

            CREATE TABLE IF NOT EXISTS members (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              project_id TEXT NOT NULL,
              name TEXT NOT NULL,
              token TEXT NOT NULL,
              UNIQUE(project_id, name),
              UNIQUE(project_id, token),
              FOREIGN KEY(project_id) REFERENCES projects(id)
            );

            -- Documents ingested from folders
            CREATE TABLE IF NOT EXISTS documents (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              project_id TEXT NOT NULL,
              dir_kind TEXT NOT NULL,      -- 'background' or 'requirement'
              filename TEXT NOT NULL,
              sha256 TEXT NOT NULL,
              extracted_text TEXT NOT NULL,
              created_at TEXT NOT NULL,
              UNIQUE(project_id, dir_kind, filename, sha256),
              FOREIGN KEY(project_id) REFERENCES projects(id)
            );

            -- Chunk storage (for retrieval)
            CREATE TABLE IF NOT EXISTS doc_chunks (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              project_id TEXT NOT NULL,
              document_id INTEGER NOT NULL,
              chunk_index INTEGER NOT NULL,
              content TEXT NOT NULL,
              created_at TEXT NOT NULL,
              FOREIGN KEY(project_id) REFERENCES projects(id),
              FOREIGN KEY(document_id) REFERENCES documents(id)
            );

            -- FTS for chunk retrieval (best-effort; some SQLite builds may lack FTS5)
            CREATE VIRTUAL TABLE IF NOT EXISTS doc_chunks_fts
            USING fts5(
              content,
              project_id UNINDEXED,
              chunk_id UNINDEXED,
              tokenize = 'porter'
            );

            -- Topics and votes
            CREATE TABLE IF NOT EXISTS topics (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              project_id TEXT NOT NULL,
              title TEXT NOT NULL,
              one_liner TEXT NOT NULL,
              rationale TEXT NOT NULL,
              feasibility TEXT NOT NULL,
              data_sources_json TEXT NOT NULL,
              key_questions_json TEXT NOT NULL,
              risks_json TEXT NOT NULL,
              score_json TEXT NOT NULL DEFAULT '{}',
              created_at TEXT NOT NULL,
              FOREIGN KEY(project_id) REFERENCES projects(id)
            );

            CREATE TABLE IF NOT EXISTS votes (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              project_id TEXT NOT NULL,
              member_id INTEGER NOT NULL,
              topic_id INTEGER NOT NULL,
              voted_at TEXT NOT NULL,
              UNIQUE(project_id, member_id),
              FOREIGN KEY(project_id) REFERENCES projects(id),
              FOREIGN KEY(member_id) REFERENCES members(id),
              FOREIGN KEY(topic_id) REFERENCES topics(id)
            );

            -- Tasks and signups
            CREATE TABLE IF NOT EXISTS tasks (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              project_id TEXT NOT NULL,
              title TEXT NOT NULL,
              description TEXT NOT NULL,
              due_date TEXT,
              max_people INTEGER NOT NULL,
              created_at TEXT NOT NULL,
              FOREIGN KEY(project_id) REFERENCES projects(id)
            );

            CREATE TABLE IF NOT EXISTS signups (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              project_id TEXT NOT NULL,
              task_id INTEGER NOT NULL,
              member_id INTEGER NOT NULL,
              signed_at TEXT NOT NULL,
              UNIQUE(project_id, task_id, member_id),
              FOREIGN KEY(project_id) REFERENCES projects(id),
              FOREIGN KEY(task_id) REFERENCES tasks(id),
              FOREIGN KEY(member_id) REFERENCES members(id)
            );

            -- Deliberation logs
            CREATE TABLE IF NOT EXISTS deliberations (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              project_id TEXT NOT NULL,
              kind TEXT NOT NULL,              -- 'topics' | 'tasks' | 'final'
              round_index INTEGER NOT NULL,
              actor TEXT NOT NULL,             -- agent name or 'controller'
              model TEXT NOT NULL,
              content_json TEXT NOT NULL,
              created_at TEXT NOT NULL,
              FOREIGN KEY(project_id) REFERENCES projects(id)
            );

            -- Final artifact outputs
            CREATE TABLE IF NOT EXISTS artifacts (
              id INTEGER PRIMARY KEY AUTOINCREMENT,
              project_id TEXT NOT NULL,
              kind TEXT NOT NULL,              -- 'final_markdown' | 'process_markdown'
              content TEXT NOT NULL,
              created_at TEXT NOT NULL,
              FOREIGN KEY(project_id) REFERENCES projects(id)
            );
            """
        )
        self.conn.commit()

    # ---- projects ----
    def create_project(self, name: str, settings: Dict[str, Any]) -> str:
        pid = uuid.uuid4().hex[:10]
        cur = self.conn.cursor()
        cur.execute(
            "INSERT INTO projects(id,name,created_at,settings_json) VALUES(?,?,?,?)",
            (pid, name, utc_now_iso(), json.dumps(settings, ensure_ascii=False)),
        )
        self.conn.commit()
        return pid

    def get_project(self, pid: str) -> sqlite3.Row:
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM projects WHERE id=?", (pid,))
        row = cur.fetchone()
        if not row:
            raise KeyError(f"Project not found: {pid}")
        return row

    def update_project(self, pid: str, **fields: Any) -> None:
        if not fields:
            return
        cols, vals = zip(*fields.items())
        sets = ", ".join([f"{c}=?" for c in cols])
        cur = self.conn.cursor()
        cur.execute(f"UPDATE projects SET {sets} WHERE id=?", (*vals, pid))
        self.conn.commit()

    # ---- members ----
    def add_members(self, pid: str, names: List[str]) -> List[Tuple[str, str]]:
        out: List[Tuple[str, str]] = []
        cur = self.conn.cursor()
        for name in names:
            name = name.strip()
            if not name:
                continue
            token = uuid.uuid4().hex
            cur.execute(
                "INSERT OR IGNORE INTO members(project_id,name,token) VALUES(?,?,?)",
                (pid, name, token),
            )
            out.append((name, token))
        self.conn.commit()
        return out

    def list_members(self, pid: str) -> List[sqlite3.Row]:
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM members WHERE project_id=? ORDER BY id ASC", (pid,))
        return list(cur.fetchall())

    def member_by_token(self, pid: str, token: str) -> Optional[sqlite3.Row]:
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM members WHERE project_id=? AND token=?", (pid, token))
        return cur.fetchone()

    # ---- documents & chunks ----
    def clear_documents(self, pid: str, dir_kind: Optional[str] = None) -> None:
        cur = self.conn.cursor()
        if dir_kind:
            cur.execute("SELECT id FROM documents WHERE project_id=? AND dir_kind=?", (pid, dir_kind))
        else:
            cur.execute("SELECT id FROM documents WHERE project_id=?", (pid,))
        doc_ids = [int(r["id"]) for r in cur.fetchall()]
        if doc_ids:
            placeholders = ",".join(["?"] * len(doc_ids))
            cur.execute(f"DELETE FROM doc_chunks WHERE document_id IN ({placeholders})", doc_ids)
            cur.execute(f"DELETE FROM doc_chunks_fts WHERE chunk_id IN (SELECT id FROM doc_chunks WHERE document_id IN ({placeholders}))", doc_ids)
        if dir_kind:
            cur.execute("DELETE FROM documents WHERE project_id=? AND dir_kind=?", (pid, dir_kind))
        else:
            cur.execute("DELETE FROM documents WHERE project_id=?", (pid,))
        self.conn.commit()

    def insert_document(self, pid: str, dir_kind: str, filename: str, extracted_text: str) -> int:
        sha = sha256_hex(extracted_text)
        cur = self.conn.cursor()
        cur.execute(
            """
            INSERT INTO documents(project_id,dir_kind,filename,sha256,extracted_text,created_at)
            VALUES(?,?,?,?,?,?)
            """,
            (pid, dir_kind, filename, sha, extracted_text, utc_now_iso()),
        )
        self.conn.commit()
        return int(cur.lastrowid)

    def list_documents(self, pid: str, dir_kind: Optional[str] = None) -> List[sqlite3.Row]:
        cur = self.conn.cursor()
        if dir_kind:
            cur.execute(
                "SELECT * FROM documents WHERE project_id=? AND dir_kind=? ORDER BY id ASC",
                (pid, dir_kind),
            )
        else:
            cur.execute("SELECT * FROM documents WHERE project_id=? ORDER BY id ASC", (pid,))
        return list(cur.fetchall())

    def insert_chunks(self, pid: str, document_id: int, chunks: List[str]) -> None:
        cur = self.conn.cursor()
        for idx, content in enumerate(chunks):
            cur.execute(
                """
                INSERT INTO doc_chunks(project_id,document_id,chunk_index,content,created_at)
                VALUES(?,?,?,?,?)
                """,
                (pid, document_id, idx, content, utc_now_iso()),
            )
            chunk_id = int(cur.lastrowid)
            # FTS insert (may fail if FTS not supported); ignore errors safely
            try:
                cur.execute(
                    "INSERT INTO doc_chunks_fts(content,project_id,chunk_id) VALUES(?,?,?)",
                    (content, pid, chunk_id),
                )
            except Exception:
                pass
        self.conn.commit()

    def search_chunks(self, pid: str, query: str, limit: int = 8) -> List[sqlite3.Row]:
        """
        FTS5 search best-effort. If FTS fails, fallback to naive LIKE scan.
        Returns rows: chunk_id, content.
        """
        cur = self.conn.cursor()
        q = query.strip()
        if not q:
            return []

        # Attempt FTS5 match
        try:
            cur.execute(
                """
                SELECT chunk_id, content
                FROM doc_chunks_fts
                WHERE doc_chunks_fts MATCH ?
                LIMIT ?
                """,
                (q, int(limit)),
            )
            rows = cur.fetchall()
            if rows:
                return list(rows)
        except Exception:
            pass

        # Fallback: naive LIKE scan (slower)
        cur.execute(
            """
            SELECT id AS chunk_id, content
            FROM doc_chunks
            WHERE project_id=? AND content LIKE ?
            LIMIT ?
            """,
            (pid, f"%{q}%", int(limit)),
        )
        return list(cur.fetchall())

    def chunk_sources(self, chunk_ids: List[int]) -> Dict[int, Dict[str, Any]]:
        """
        Map chunk_id -> {filename, dir_kind, chunk_index}.
        """
        if not chunk_ids:
            return {}
        cur = self.conn.cursor()
        placeholders = ",".join(["?"] * len(chunk_ids))
        cur.execute(
            f"""
            SELECT c.id AS chunk_id, c.chunk_index, d.filename, d.dir_kind
            FROM doc_chunks c
            JOIN documents d ON d.id=c.document_id
            WHERE c.id IN ({placeholders})
            """,
            chunk_ids,
        )
        out: Dict[int, Dict[str, Any]] = {}
        for r in cur.fetchall():
            out[int(r["chunk_id"])] = {
                "filename": r["filename"],
                "dir_kind": r["dir_kind"],
                "chunk_index": int(r["chunk_index"]),
            }
        return out

    # ---- topics ----
    def clear_topics(self, pid: str) -> None:
        cur = self.conn.cursor()
        cur.execute("DELETE FROM topics WHERE project_id=?", (pid,))
        cur.execute("DELETE FROM votes WHERE project_id=?", (pid,))
        self.conn.commit()

    def insert_topics(self, pid: str, topics: List[Dict[str, Any]]) -> None:
        cur = self.conn.cursor()
        for t in topics:
            cur.execute(
                """
                INSERT INTO topics(
                  project_id,title,one_liner,rationale,feasibility,
                  data_sources_json,key_questions_json,risks_json,score_json,created_at
                ) VALUES(?,?,?,?,?,?,?,?,?,?)
                """,
                (
                    pid,
                    t["title"],
                    t["one_liner"],
                    t["rationale"],
                    t["feasibility"],
                    json.dumps(t.get("data_sources", []), ensure_ascii=False),
                    json.dumps(t.get("key_questions", []), ensure_ascii=False),
                    json.dumps(t.get("risks", []), ensure_ascii=False),
                    json.dumps(t.get("scores", {}), ensure_ascii=False),
                    utc_now_iso(),
                ),
            )
        self.conn.commit()

    def list_topics(self, pid: str) -> List[sqlite3.Row]:
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM topics WHERE project_id=? ORDER BY id ASC", (pid,))
        return list(cur.fetchall())

    def topic_by_id(self, pid: str, topic_id: int) -> Optional[sqlite3.Row]:
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM topics WHERE project_id=? AND id=?", (pid, topic_id))
        return cur.fetchone()

    # ---- voting ----
    def cast_vote(self, pid: str, member_id: int, topic_id: int) -> None:
        cur = self.conn.cursor()
        cur.execute(
            "INSERT OR REPLACE INTO votes(project_id,member_id,topic_id,voted_at) VALUES(?,?,?,?)",
            (pid, member_id, topic_id, utc_now_iso()),
        )
        self.conn.commit()

    def tally_votes(self, pid: str) -> List[sqlite3.Row]:
        cur = self.conn.cursor()
        cur.execute(
            """
            SELECT t.id AS topic_id, t.title AS title, COUNT(v.id) AS votes
            FROM topics t
            LEFT JOIN votes v ON t.id=v.topic_id AND v.project_id=?
            WHERE t.project_id=?
            GROUP BY t.id
            ORDER BY votes DESC, t.id ASC
            """,
            (pid, pid),
        )
        return list(cur.fetchall())

    def selected_topic(self, pid: str) -> Optional[sqlite3.Row]:
        rows = self.tally_votes(pid)
        if not rows:
            return None
        top = rows[0]
        return self.topic_by_id(pid, int(top["topic_id"]))

    # ---- tasks ----
    def clear_tasks(self, pid: str) -> None:
        cur = self.conn.cursor()
        cur.execute("DELETE FROM tasks WHERE project_id=?", (pid,))
        cur.execute("DELETE FROM signups WHERE project_id=?", (pid,))
        self.conn.commit()

    def insert_tasks(self, pid: str, tasks: List[Dict[str, Any]]) -> None:
        cur = self.conn.cursor()
        for t in tasks:
            cur.execute(
                """
                INSERT INTO tasks(project_id,title,description,due_date,max_people,created_at)
                VALUES(?,?,?,?,?,?)
                """,
                (
                    pid,
                    t["title"],
                    t["description"],
                    t.get("due_date"),
                    int(t.get("max_people", 1)),
                    utc_now_iso(),
                ),
            )
        self.conn.commit()

    def list_tasks(self, pid: str) -> List[sqlite3.Row]:
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM tasks WHERE project_id=? ORDER BY id ASC", (pid,))
        return list(cur.fetchall())

    def signup_counts(self, pid: str) -> Dict[int, int]:
        cur = self.conn.cursor()
        cur.execute("SELECT task_id, COUNT(*) AS c FROM signups WHERE project_id=? GROUP BY task_id", (pid,))
        return {int(r["task_id"]): int(r["c"]) for r in cur.fetchall()}

    def list_signups_for_task(self, pid: str, task_id: int) -> List[sqlite3.Row]:
        cur = self.conn.cursor()
        cur.execute(
            """
            SELECT s.*, m.name AS member_name
            FROM signups s
            JOIN members m ON m.id=s.member_id
            WHERE s.project_id=? AND s.task_id=?
            ORDER BY s.signed_at ASC
            """,
            (pid, task_id),
        )
        return list(cur.fetchall())

    def member_current_signup(self, pid: str, member_id: int) -> Optional[int]:
        cur = self.conn.cursor()
        cur.execute("SELECT task_id FROM signups WHERE project_id=? AND member_id=? LIMIT 1", (pid, member_id))
        row = cur.fetchone()
        if not row:
            return None
        return int(row["task_id"])

    def signup_task(self, pid: str, member_id: int, task_id: int, allow_change: bool = True) -> Tuple[bool, str]:
        cur = self.conn.cursor()
        cur.execute("SELECT * FROM tasks WHERE project_id=? AND id=?", (pid, task_id))
        task = cur.fetchone()
        if not task:
            return False, "Task not found."

        existing = self.member_current_signup(pid, member_id)
        if existing is not None and not allow_change:
            return False, "You already signed up for a task. (Changing disabled)"
        if existing is not None and allow_change and existing != task_id:
            cur.execute("DELETE FROM signups WHERE project_id=? AND member_id=?", (pid, member_id))

        cur.execute("SELECT COUNT(*) AS c FROM signups WHERE project_id=? AND task_id=?", (pid, task_id))
        c = int(cur.fetchone()["c"])
        if c >= int(task["max_people"]):
            return False, "This task is full."

        cur.execute(
            "INSERT OR IGNORE INTO signups(project_id,task_id,member_id,signed_at) VALUES(?,?,?,?)",
            (pid, task_id, member_id, utc_now_iso()),
        )
        self.conn.commit()
        return True, "Signed up!"

    # ---- deliberations ----
    def log_deliberation(self, pid: str, kind: str, round_index: int, actor: str, model: str, content: Dict[str, Any]) -> None:
        cur = self.conn.cursor()
        cur.execute(
            """
            INSERT INTO deliberations(project_id,kind,round_index,actor,model,content_json,created_at)
            VALUES(?,?,?,?,?,?,?)
            """,
            (pid, kind, round_index, actor, model, json.dumps(content, ensure_ascii=False), utc_now_iso()),
        )
        self.conn.commit()

    def list_deliberations(self, pid: str, kind: Optional[str] = None) -> List[sqlite3.Row]:
        cur = self.conn.cursor()
        if kind:
            cur.execute(
                "SELECT * FROM deliberations WHERE project_id=? AND kind=? ORDER BY id ASC",
                (pid, kind),
            )
        else:
            cur.execute("SELECT * FROM deliberations WHERE project_id=? ORDER BY id ASC", (pid,))
        return list(cur.fetchall())

    # ---- artifacts ----
    def upsert_artifact(self, pid: str, kind: str, content: str) -> None:
        cur = self.conn.cursor()
        cur.execute("DELETE FROM artifacts WHERE project_id=? AND kind=?", (pid, kind))
        cur.execute(
            "INSERT INTO artifacts(project_id,kind,content,created_at) VALUES(?,?,?,?)",
            (pid, kind, content, utc_now_iso()),
        )
        self.conn.commit()

    def get_artifact(self, pid: str, kind: str) -> Optional[str]:
        cur = self.conn.cursor()
        cur.execute("SELECT content FROM artifacts WHERE project_id=? AND kind=? ORDER BY id DESC LIMIT 1", (pid, kind))
        row = cur.fetchone()
        if not row:
            return None
        return str(row["content"])


# ----------------------------- Chunking / RAG ---------------------------------

def chunk_text(text: str, max_chars: int = 1800, overlap: int = 200) -> List[str]:
    """
    Simple chunker: breaks text into ~max_chars chunks with overlaps.
    """
    text = re.sub(r"\r\n?", "\n", text).strip()
    if not text:
        return []
    chunks: List[str] = []
    i = 0
    n = len(text)
    while i < n:
        j = min(n, i + max_chars)
        chunk = text[i:j].strip()
        if chunk:
            chunks.append(chunk)
        if j >= n:
            break
        i = max(0, j - overlap)
    return chunks


def compose_rag_context(db: DB, pid: str, query: str, k: int = 8) -> Tuple[str, List[Dict[str, Any]]]:
    """
    Retrieves top-k relevant chunks and formats them with lightweight source tags.
    Returns (context_text, sources_list).
    """
    rows = db.search_chunks(pid, query=query, limit=k)
    chunk_ids = [int(r["chunk_id"]) for r in rows]
    meta = db.chunk_sources(chunk_ids)
    sources: List[Dict[str, Any]] = []
    blocks: List[str] = []
    for r in rows:
        cid = int(r["chunk_id"])
        src = meta.get(cid, {"filename": "unknown", "dir_kind": "unknown", "chunk_index": -1})
        tag = f"[{src['dir_kind']}:{src['filename']}#chunk{src['chunk_index']}]"
        content = str(r["content"]).strip()
        blocks.append(f"{tag}\n{content}")
        sources.append({"chunk_id": cid, **src})
    if not blocks:
        return "(no relevant chunks found)", []
    return "\n\n---\n\n".join(blocks), sources


# ----------------------------- OpenAI-compatible client ------------------------

@dataclass
class LLMCallConfig:
    api_base: str
    api_key: str
    temperature: float = 0.4
    max_tokens: int = 2800
    # OpenRouter recommends these optional headers; keep customizable.
    extra_headers: Dict[str, str] = None  # type: ignore


class LLMClient:
    def __init__(self, cfg: LLMCallConfig):
        self.cfg = cfg
        if self.cfg.extra_headers is None:
            self.cfg.extra_headers = {}

    def chat(
        self,
        model: str,
        messages: List[Dict[str, str]],
        *,
        json_mode: bool = False,
        max_tokens: Optional[int] = None,
        temperature: Optional[float] = None,
        timeout: int = 180,
    ) -> str:
        url = self.cfg.api_base.rstrip("/") + "/chat/completions"
        headers = {
            "Authorization": f"Bearer {self.cfg.api_key}",
            "Content-Type": "application/json",
            **(self.cfg.extra_headers or {}),
        }
        body: Dict[str, Any] = {
            "model": model,
            "messages": messages,
            "temperature": float(self.cfg.temperature if temperature is None else temperature),
            "max_tokens": int(self.cfg.max_tokens if max_tokens is None else max_tokens),
        }
        if json_mode:
            body["response_format"] = {"type": "json_object"}

        resp = requests.post(url, headers=headers, json=body, timeout=timeout)
        if resp.status_code != 200:
            raise RuntimeError(f"LLM API error HTTP {resp.status_code}: {resp.text[:1500]}")
        data = resp.json()
        try:
            content = data["choices"][0]["message"]["content"]
        except Exception:
            raise RuntimeError(f"Unexpected response shape: {data}")
        if isinstance(content, list):
            content = "".join(part.get("text", "") if isinstance(part, dict) else str(part) for part in content)
        return str(content)

    def chat_json(
        self,
        model: str,
        messages: List[Dict[str, str]],
        *,
        retries: int = 2,
        timeout: int = 180,
    ) -> Dict[str, Any]:
        """
        JSON-only contract with auto-repair.
        """
        raw = self.chat(model=model, messages=messages, json_mode=True, timeout=timeout)
        for attempt in range(retries + 1):
            try:
                return try_parse_json(raw)
            except Exception:
                if attempt >= retries:
                    raise RuntimeError(f"Model returned invalid JSON after retries. Raw:\n{raw[:2000]}")
                repair = [
                    {"role": "system", "content": "You are a strict JSON repair bot. Output ONLY valid JSON object. No extra text."},
                    {"role": "user", "content": f"Fix the following into a valid JSON object, preserving keys/values as much as possible:\n\n{raw}"},
                ]
                raw = self.chat(model=model, messages=repair, json_mode=True, timeout=timeout)
        raise RuntimeError("unreachable")


# ----------------------------- Agents & schemas --------------------------------

@dataclass
class Agent:
    name: str
    model: str
    role_prompt: str


def default_agents(models: Dict[str, str]) -> List[Agent]:
    """
    "Extreme" roster. You can change/extend in project settings or CLI flags.
    """
    return [
        Agent(
            name="Leader",
            model=models["leader"],
            role_prompt=(
                "You are the group leader. You merge conflicting suggestions into a coherent plan. "
                "You care about structure, clarity, rubric alignment, and end-to-end feasibility."
            ),
        ),
        Agent(
            name="Critic",
            model=models["critic"],
            role_prompt=(
                "You are the hard critic. You find holes, weak logic, missing sections, and mismatched deliverables. "
                "You propose fixes with concrete examples."
            ),
        ),
        Agent(
            name="Researcher",
            model=models["researcher"],
            role_prompt=(
                "You are the researcher. You focus on concrete methods, data sources, realistic scope, and how to produce evidence."
            ),
        ),
        Agent(
            name="Methodologist",
            model=models["methodologist"],
            role_prompt=(
                "You are the methodologist. You design study methods (survey/interview/case analysis/metrics), validity threats, sampling, ethics."
            ),
        ),
        Agent(
            name="Editor",
            model=models["editor"],
            role_prompt=(
                "You are the editor. You rewrite for readability, consistent headings, strong executive summary, crisp bullets, and clean markdown."
            ),
        ),
        Agent(
            name="RedTeamer",
            model=models["redteam"],
            role_prompt=(
                "You are the red-teamer. You stress-test feasibility, check for hallucination risks, and demand that claims be grounded in provided docs."
            ),
        ),
    ]


TOPIC_SCHEMA_HINT = {
    "topics": [
        {
            "title": "string",
            "one_liner": "string",
            "rationale": "string",
            "feasibility": "string",
            "data_sources": ["string", "..."],
            "key_questions": ["string", "..."],
            "risks": ["string", "..."],
            "scores": {
                "rubric_fit": 0,
                "feasibility": 0,
                "originality": 0,
                "data_availability": 0,
                "presentation_value": 0
            }
        }
    ]
}

TASK_SCHEMA_HINT = {
    "tasks": [
        {
            "title": "string",
            "description": "string (must include deliverables + acceptance criteria)",
            "max_people": 1,
            "due_date": "YYYY-MM-DD or null"
        }
    ]
}


# ----------------------------- Deliberation logic ------------------------------

def _controller_should_continue_prompt() -> str:
    return textwrap.dedent("""
    You are the CoT Controller.
    You may reason internally step-by-step, but you MUST NOT reveal chain-of-thought.
    Output ONLY a JSON object of the form:
    {
      "should_continue": true/false,
      "quality_score": number 0..10,
      "reason": "1-3 sentences",
      "next_focus": ["short bullet", "..."]   // optional but recommended
    }
    """).strip()


def deliberate_topics(
    db: DB,
    llm: LLMClient,
    pid: str,
    agents: List[Agent],
    controller_model: str,
    *,
    n_topics: int = 8,
    min_rounds: int = 2,
    max_rounds: int = 6,
    rag_k: int = 8,
) -> List[Dict[str, Any]]:
    """
    Multi-round topic deliberation:
      - Each agent proposes/refines topics in JSON.
      - Leader merges/dedupes.
      - Controller decides continue/stop.
    Stores logs in DB and returns final topics list.
    """
    n_topics = clamp(n_topics, 3, 12)
    min_rounds = clamp(min_rounds, 1, 12)
    max_rounds = clamp(max_rounds, min_rounds, 20)

    seed_query = "grading rubric deliverables requirements topic constraints"
    rag_context, sources = compose_rag_context(db, pid, seed_query, k=rag_k)
    src_note = "\n".join([f"- [{s['dir_kind']}:{s['filename']}#chunk{s['chunk_index']}]" for s in sources]) or "(none)"

    leader = agents[0]

    # Round state
    current_topics: List[Dict[str, Any]] = []
    round_index = 0

    def agent_prompt(agent: Agent, round_i: int, current: List[Dict[str, Any]], focus: List[str]) -> List[Dict[str, str]]:
        focus_text = "\n".join([f"- {x}" for x in focus]) if focus else "(no special focus)"
        current_json = json.dumps({"topics": current}, ensure_ascii=False, indent=2) if current else "(none yet)"
        return [
            {
                "role": "system",
                "content": textwrap.dedent(f"""
                You are {agent.name}.
                Role: {agent.role_prompt}

                Task: Propose or refine {n_topics} candidate research topics for the assignment,
                grounded in provided documents. Avoid making up stats or fake citations.
                Output ONLY JSON with key "topics" following this example shape:
                {json.dumps(TOPIC_SCHEMA_HINT, ensure_ascii=False)}
                """).strip(),
            },
            {
                "role": "user",
                "content": textwrap.dedent(f"""
                RAG EXCERPTS (most relevant chunks):
                ---
                {rag_context}
                ---

                Source tags present in excerpts:
                {src_note}

                Current round: {round_i}
                Current topic set (JSON):
                {current_json}

                Controller focus for this round:
                {focus_text}

                Requirements:
                - Return EXACTLY {n_topics} topics.
                - Each topic must include: title, one_liner, rationale, feasibility, data_sources, key_questions, risks, scores.
                - Scores are 0..10 integers (rough rubric-based estimation).
                - Keep topics feasible within typical course time constraints (1-3 weeks unless docs say otherwise).
                """).strip(),
            },
        ]

    def leader_merge_prompt(
        round_i: int,
        proposals: List[Tuple[Agent, Dict[str, Any]]],
        focus: List[str]
    ) -> List[Dict[str, str]]:

        merged = []
        for ag, d in proposals:
            merged.append(
                f"=== {ag.name} ({ag.model}) ===\n"
                f"{json.dumps(d, ensure_ascii=False, indent=2)}\n"
            )

        # ✅ 所有中间变量，提前算好
        merged_text = "\n\n".join(merged)
        focus_text = "\n".join([f"- {x}" for x in focus]) if focus else "(none)"

        return [
            {
                "role": "system",
                "content": textwrap.dedent(f"""
                You are the Leader and synthesizer.
                Merge multiple topic lists into a single high-quality set of EXACTLY {n_topics} topics.

                Output ONLY JSON with:
                {{
                "topics": [...],
                "merge_notes": ["short note", "..."],
                "dedupe_map": {{"old_title":"new_title_or_same", "..."}}
                }}

                Rules:
                - Deduplicate near-duplicates.
                - Ensure coverage across plausible angles suggested by docs.
                - Ensure each topic has the required fields and scores 0..10.
                - Avoid hallucinated citations; use only doc-based phrasing.
                """).strip(),
            },
            {
                "role": "user",
                "content": textwrap.dedent(f"""
                Current round: {round_i}
                Controller focus:
                {focus_text}

                Here are the proposals:
                {merged_text}
                """).strip(),
            },
        ]

    def controller_prompt(round_i: int, topics: List[Dict[str, Any]], focus_prev: List[str]) -> List[Dict[str, str]]:
        topics_json = json.dumps({"topics": topics}, ensure_ascii=False, indent=2)
        return [
            {"role": "system", "content": _controller_should_continue_prompt()},
            {
                "role": "user",
                "content": textwrap.dedent(f"""
                Evaluate the current topic set quality and decide if more rounds are needed.

                RAG context sources:
                {src_note}

                Round: {round_i}
                Previous focus:
                {'; '.join(focus_prev) if focus_prev else '(none)'}

                Current topics:
                {topics_json}

                Consider:
                - rubric fit
                - feasibility
                - data availability
                - diversity of angles
                - clarity of titles/one-liners
                """).strip(),
            },
        ]

    focus: List[str] = [
        "Ensure topics map cleanly to deliverables and grading rubric.",
        "Ensure each topic has realistic data sources the student can access quickly.",
        "Avoid vague topics; make them operationalizable.",
    ]

    while True:
        round_index += 1

        # Each agent proposes/refines
        proposals: List[Tuple[Agent, Dict[str, Any]]] = []
        for ag in agents:
            data = llm.chat_json(ag.model, agent_prompt(ag, round_index, current_topics, focus))
            db.log_deliberation(pid, "topics", round_index, ag.name, ag.model, data)
            proposals.append((ag, data))

        # Leader merges
        merged = llm.chat_json(leader.model, leader_merge_prompt(round_index, proposals, focus))
        db.log_deliberation(pid, "topics", round_index, "LeaderMerge", leader.model, merged)

        topics = merged.get("topics", [])
        if not isinstance(topics, list) or len(topics) != n_topics:
            # Hard fallback: accept first n_topics from whatever we got
            topics = topics if isinstance(topics, list) else []
            topics = topics[:n_topics]
            while len(topics) < n_topics:
                topics.append({
                    "title": f"Topic Placeholder {len(topics)+1}",
                    "one_liner": "Placeholder due to invalid merge output.",
                    "rationale": "N/A",
                    "feasibility": "N/A",
                    "data_sources": [],
                    "key_questions": [],
                    "risks": [],
                    "scores": {"rubric_fit": 0, "feasibility": 0, "originality": 0, "data_availability": 0, "presentation_value": 0},
                })

        current_topics = topics

        # Controller decides continue
        decision = llm.chat_json(controller_model, controller_prompt(round_index, current_topics, focus))
        db.log_deliberation(pid, "topics", round_index, "Controller", controller_model, decision)

        should_continue = bool(decision.get("should_continue", False))
        quality = float(decision.get("quality_score", 0))
        next_focus = decision.get("next_focus", [])
        if not isinstance(next_focus, list):
            next_focus = []

        # Update focus adaptively (extreme logic: keep pushing on weak dimensions)
        focus = []
        if quality < 7.5:
            focus.append("Improve rubric alignment and ensure deliverables are obvious.")
            focus.append("Increase feasibility: tighter scope, clear methods, clear data collection plan.")
        if quality < 8.5:
            focus.append("Improve topic differentiation and reduce redundancy.")
            focus.append("Strengthen key questions and risks.")
        focus += [str(x) for x in next_focus][:5]

        if round_index < min_rounds:
            continue
        if (not should_continue) or (round_index >= max_rounds):
            break

    return current_topics


def auto_pick_best_topic(
    llm: LLMClient,
    model: str,
    topics: List[Dict[str, Any]],
) -> Dict[str, Any]:
    """
    Ask a model to pick the best topic with a short explanation + tie-break rules.
    """
    prompt = [
        {"role": "system", "content": "Pick the single best topic. Output ONLY JSON: {\"selected_title\":...,\"reason\":...,\"tie_break\":...}"},
        {"role": "user", "content": f"Topics JSON:\n{json.dumps({'topics': topics}, ensure_ascii=False, indent=2)}\n\nPick the best overall."},
    ]
    return llm.chat_json(model, prompt)


def deliberate_tasks(
    db: DB,
    llm: LLMClient,
    pid: str,
    agents: List[Agent],
    controller_model: str,
    *,
    topic_title: str,
    member_count: int,
    min_rounds: int = 2,
    max_rounds: int = 6,
    rag_k: int = 8,
) -> List[Dict[str, Any]]:
    """
    Multi-round task breakdown deliberation:
      - Each agent proposes tasks JSON
      - Leader merges into a robust task list
      - Controller decides to continue/stop
    """
    member_count = max(1, member_count)
    # "extreme" default: tasks >= members, but not too many
    target_tasks = clamp(member_count + 3, 6, 24)

    rag_context, sources = compose_rag_context(db, pid, f"deliverables sections format requirements {topic_title}", k=rag_k)
    src_note = "\n".join([f"- [{s['dir_kind']}:{s['filename']}#chunk{s['chunk_index']}]" for s in sources]) or "(none)"
    leader = agents[0]
    round_index = 0
    current_tasks: List[Dict[str, Any]] = []

    focus: List[str] = [
        "Tasks must be parallelizable with minimal dependencies.",
        "Each task must specify deliverables + acceptance criteria.",
        "Include a timeline and integration checkpoints.",
        "Cover slides/report + data/method + QA + references.",
    ]

    def agent_prompt(agent: Agent, round_i: int, current: List[Dict[str, Any]], focus_list: List[str]) -> List[Dict[str, str]]:
        focus_text = "\n".join([f"- {x}" for x in focus_list]) if focus_list else "(none)"
        cur_json = json.dumps({"tasks": current}, ensure_ascii=False, indent=2) if current else "(none yet)"
        return [
            {
                "role": "system",
                "content": textwrap.dedent(f"""
                You are {agent.name}.
                Role: {agent.role_prompt}

                Task: Propose or refine a task list for topic: "{topic_title}".
                Output ONLY JSON with key "tasks", following this example shape:
                {json.dumps(TASK_SCHEMA_HINT, ensure_ascii=False)}

                Hard rules:
                - Return EXACTLY {target_tasks} tasks.
                - Each task must include deliverables + acceptance criteria in description.
                - Use due_date as YYYY-MM-DD if you can infer schedule, else null.
                - max_people between 1 and 6.
                """).strip(),
            },
            {
                "role": "user",
                "content": textwrap.dedent(f"""
                RAG EXCERPTS:
                ---
                {rag_context}
                ---

                Source tags:
                {src_note}

                Round: {round_i}
                Current tasks (JSON):
                {cur_json}

                Focus:
                {focus_text}

                Constraints:
                - Assume a small team equivalent; but in reality a single student can do tasks sequentially with AI help.
                - Still keep tasks realistic and not too tiny.
                """).strip(),
            },
        ]

    def leader_merge_prompt(
        round_i: int,
        proposals: List[Tuple[Agent, Dict[str, Any]]],
        focus_list: List[str],
    ) -> List[Dict[str, str]]:
        merged = []
        for ag, d in proposals:
            merged.append(
                f"=== {ag.name} ({ag.model}) ===\n"
                f"{json.dumps(d, ensure_ascii=False, indent=2)}\n"
            )

        focus_text = "\n".join([f"- {x}" for x in focus_list]) if focus_list else "(none)"

        # ✅ precompute (avoid backslash inside f-string expressions)
        proposals_text = "\n\n".join(merged)

        return [
            {
                "role": "system",
                "content": textwrap.dedent(f"""
                You are the Leader synthesizer.
                Merge multiple task lists into a single robust list of EXACTLY {target_tasks} tasks.

                Output ONLY JSON:
                {{
                "tasks": [...],
                "merge_notes": ["..."],
                "dependency_graph": [
                    {{"task":"A","depends_on":["B","C"]}},
                    ...
                ],
                "integration_plan": ["checkpoint 1", "checkpoint 2", ...]
                }}

                Rules:
                - Remove duplicates, ensure full coverage.
                - Tasks must include deliverables + acceptance criteria.
                - Keep a reasonable ordering and note dependencies.
                """).strip(),
            },
            {
                "role": "user",
                "content": textwrap.dedent(f"""
                Round: {round_i}
                Focus:
                {focus_text}

                Proposals:
                {proposals_text}
                """).strip(),
            },
        ]

    def controller_prompt(round_i: int, tasks: List[Dict[str, Any]], focus_prev: List[str]) -> List[Dict[str, str]]:
        return [
            {"role": "system", "content": _controller_should_continue_prompt()},
            {
                "role": "user",
                "content": textwrap.dedent(f"""
                Decide whether the task list is complete and robust enough.

                Sources:
                {src_note}

                Round: {round_i}
                Previous focus:
                {'; '.join(focus_prev) if focus_prev else '(none)'}

                Task list:
                {json.dumps({'tasks': tasks}, ensure_ascii=False, indent=2)}

                Consider:
                - coverage of deliverables
                - realism
                - parallelism and dependencies
                - acceptance criteria quality
                """).strip(),
            },
        ]

    while True:
        round_index += 1
        proposals: List[Tuple[Agent, Dict[str, Any]]] = []
        for ag in agents:
            data = llm.chat_json(ag.model, agent_prompt(ag, round_index, current_tasks, focus))
            db.log_deliberation(pid, "tasks", round_index, ag.name, ag.model, data)
            proposals.append((ag, data))

        merged = llm.chat_json(leader.model, leader_merge_prompt(round_index, proposals, focus))
        db.log_deliberation(pid, "tasks", round_index, "LeaderMerge", leader.model, merged)

        tasks = merged.get("tasks", [])
        if not isinstance(tasks, list) or len(tasks) != target_tasks:
            tasks = tasks if isinstance(tasks, list) else []
            tasks = tasks[:target_tasks]
            while len(tasks) < target_tasks:
                tasks.append({"title": f"Task Placeholder {len(tasks)+1}", "description": "Placeholder", "max_people": 1, "due_date": None})
        current_tasks = tasks

        decision = llm.chat_json(controller_model, controller_prompt(round_index, current_tasks, focus))
        db.log_deliberation(pid, "tasks", round_index, "Controller", controller_model, decision)

        quality = float(decision.get("quality_score", 0))
        should_continue = bool(decision.get("should_continue", False))
        next_focus = decision.get("next_focus", [])
        if not isinstance(next_focus, list):
            next_focus = []

        focus = []
        if quality < 7.5:
            focus.append("Improve acceptance criteria and deliverables clarity.")
            focus.append("Fill missing coverage: references, QA, synthesis, slide deck, appendix.")
        if quality < 8.5:
            focus.append("Improve dependency graph and integration checkpoints.")
        focus += [str(x) for x in next_focus][:5]

        if round_index < min_rounds:
            continue
        if (not should_continue) or (round_index >= max_rounds):
            break

    return current_tasks


def deliberate_final_markdown(
    db: DB,
    llm: LLMClient,
    pid: str,
    agents: List[Agent],
    controller_model: str,
    *,
    selected_topic: Dict[str, Any],
    tasks: List[Dict[str, Any]],
    min_rounds: int = 2,
    max_rounds: int = 6,
    rag_k: int = 10,
) -> str:
    """
    Multi-round final Markdown deliberation:
      - Agents propose full Markdown plan/report skeleton (not full essay, but very complete)
      - Leader merges
      - Controller decides continue
    Stores final artifact in DB.
    """
    leader = agents[0]
    round_index = 0
    current_md = ""

    rag_context, sources = compose_rag_context(db, pid, f"{selected_topic.get('title','')} deliverables format rubric", k=rag_k)
    src_note = "\n".join([f"- [{s['dir_kind']}:{s['filename']}#chunk{s['chunk_index']}]" for s in sources]) or "(none)"

    focus: List[str] = [
        "Produce a comprehensive Markdown deliverable aligned with requirements.",
        "Use the provided docs; include inline [background:...#chunk] style tags only if helpful.",
        "Include methodology, timeline, risk register, and slide outline.",
        "Avoid fabricated statistics; mark placeholders clearly.",
    ]

    def agent_prompt(agent: Agent, round_i: int, current: str, focus_list: List[str]) -> List[Dict[str, str]]:
        focus_text = "\n".join([f"- {x}" for x in focus_list]) if focus_list else "(none)"
        cur = current.strip() or "(none yet)"
        return [
            {
                "role": "system",
                "content": textwrap.dedent(f"""
                You are {agent.name}.
                Role: {agent.role_prompt}

                Task: Draft an extremely complete Markdown deliverable (report plan + near-final content structure)
                for the selected topic. Keep it feasible and grounded in provided docs. Avoid inventing facts.
                Output ONLY JSON:
                {{
                  "critic": "short critique of current draft (if any)",
                  "markdown": "your full improved Markdown draft"
                }}
                """).strip(),
            },
            {
                "role": "user",
                "content": textwrap.dedent(f"""
                RAG EXCERPTS:
                ---
                {rag_context}
                ---

                Sources:
                {src_note}

                Selected Topic:
                {json.dumps(selected_topic, ensure_ascii=False, indent=2)}

                Task list:
                {json.dumps({'tasks': tasks}, ensure_ascii=False, indent=2)}

                Round: {round_i}
                Current shared draft:
                >>> BEGIN
                {cur}
                >>> END

                Focus:
                {focus_text}

                Requirements:
                - The Markdown must include:
                  * Title + Executive summary
                  * Problem statement + objectives
                  * Research questions / hypotheses (if applicable)
                  * Methodology (data sources, design, sampling, analysis)
                  * Work plan & timeline
                  * Task responsibilities mapping (even if single student executes sequentially)
                  * Risks & mitigations
                  * Expected results & how to present (charts/tables)
                  * References strategy (what to cite)
                  * Slide deck outline (PPT structure)
                - Use placeholders where needed, clearly labeled.
                """).strip(),
            },
        ]

    def leader_merge_prompt(
        round_i: int,
        proposals: List[Tuple[Agent, Dict[str, Any]]],
        focus_list: List[str],
    ) -> List[Dict[str, str]]:
        merged = []
        for ag, d in proposals:
            merged.append(
                f"=== {ag.name} ({ag.model}) ===\n"
                f"{json.dumps(d, ensure_ascii=False, indent=2)}\n"
            )

        focus_text = "\n".join([f"- {x}" for x in focus_list]) if focus_list else "(none)"

        # ✅ precompute to avoid backslash inside f-string expression
        drafts_text = "\n\n".join(merged)

        return [
            {
                "role": "system",
                "content": textwrap.dedent("""
                You are the Leader synthesizer.
                Merge all proposed drafts into ONE best Markdown draft.
                Output ONLY JSON: {"markdown":"...","merge_notes":["..."]}

                Rules:
                - Keep a clean, consistent structure.
                - Remove duplication.
                - Strengthen weak sections and add missing sections.
                - Keep it realistic, grounded, and avoid invented facts.
                """).strip(),
            },
            {
                "role": "user",
                "content": textwrap.dedent(f"""
                Round: {round_i}
                Focus:
                {focus_text}

                Drafts:
                {drafts_text}
                """).strip(),
            },
        ]

    def controller_prompt(round_i: int, prev_md: str, new_md: str, focus_prev: List[str]) -> List[Dict[str, str]]:
        return [
            {"role": "system", "content": _controller_should_continue_prompt()},
            {
                "role": "user",
                "content": textwrap.dedent(f"""
                Decide whether the new Markdown draft is sufficiently complete.

                Sources:
                {src_note}

                Round: {round_i}
                Previous focus:
                {'; '.join(focus_prev) if focus_prev else '(none)'}

                Previous draft:
                >>> BEGIN PREV
                {prev_md[:6000]}
                >>> END PREV

                New draft:
                >>> BEGIN NEW
                {new_md[:6000]}
                >>> END NEW

                Evaluate:
                - completeness vs requirements
                - clarity and structure
                - feasibility
                - groundedness (no hallucinated facts)
                """).strip(),
            },
        ]

    while True:
        round_index += 1
        prev = current_md
        proposals: List[Tuple[Agent, Dict[str, Any]]] = []
        for ag in agents:
            d = llm.chat_json(ag.model, agent_prompt(ag, round_index, current_md, focus))
            db.log_deliberation(pid, "final", round_index, ag.name, ag.model, d)
            proposals.append((ag, d))

        merged = llm.chat_json(leader.model, leader_merge_prompt(round_index, proposals, focus))
        db.log_deliberation(pid, "final", round_index, "LeaderMerge", leader.model, merged)
        new_md = str(merged.get("markdown", "")).strip()
        if not new_md:
            new_md = prev or "# Draft\n\n(Empty output from merge; retry suggested.)"
        current_md = new_md

        decision = llm.chat_json(controller_model, controller_prompt(round_index, prev, current_md, focus))
        db.log_deliberation(pid, "final", round_index, "Controller", controller_model, decision)

        quality = float(decision.get("quality_score", 0))
        should_continue = bool(decision.get("should_continue", False))
        next_focus = decision.get("next_focus", [])
        if not isinstance(next_focus, list):
            next_focus = []

        focus = []
        if quality < 7.5:
            focus.append("Add missing required sections and tighten methodology.")
            focus.append("Make deliverables explicit, with placeholders for data/figures.")
        if quality < 8.5:
            focus.append("Improve readability and remove redundancy.")
            focus.append("Strengthen risks/mitigations and timeline.")
        focus += [str(x) for x in next_focus][:6]

        if round_index < min_rounds:
            continue
        if (not should_continue) or (round_index >= max_rounds):
            break

    db.upsert_artifact(pid, "final_markdown", current_md)
    return current_md


# ----------------------------- Leader selection --------------------------------

def pick_leader_ai(
    llm: LLMClient,
    model: str,
    member_names: List[str],
) -> Dict[str, Any]:
    """
    AI-based leader pick (not 'fair' like commit-reveal, but "smart choice" option).
    """
    prompt = [
        {"role": "system", "content": "Output ONLY JSON: {\"leader\":\"name\",\"reason\":\"...\",\"criteria\":[\"...\"],\"fallback\":\"name\"}"},
        {"role": "user", "content": f"Members: {member_names}\nPick the best leader for coordination reliability and clarity."},
    ]
    return llm.chat_json(model, prompt)


def pick_leader_commit_reveal(members: List[str], method: str, seed: Optional[str]) -> Tuple[str, Dict[str, Any]]:
    if not members:
        raise ValueError("No members.")
    seed = seed or uuid.uuid4().hex
    rng = random.Random(seed)
    if method == "lottery":
        leader = rng.choice(members)
        return leader, {"method": "lottery", "seed": seed, "commit_sha256": sha256_hex(seed)}
    # dice
    rolls = {m: rng.randint(1, 6) for m in members}
    maxv = max(rolls.values())
    candidates = [m for m, v in rolls.items() if v == maxv]
    rng.shuffle(candidates)
    leader = candidates[0]
    return leader, {"method": "dice", "seed": seed, "commit_sha256": sha256_hex(seed), "rolls": rolls}


# ----------------------------- Export markdown ---------------------------------

def export_process_markdown(db: DB, pid: str) -> str:
    p = db.get_project(pid)
    members = db.list_members(pid)
    docs_bg = db.list_documents(pid, "background")
    docs_req = db.list_documents(pid, "requirement")
    topics = db.list_topics(pid)
    tally = db.tally_votes(pid)
    selected = db.selected_topic(pid)
    tasks = db.list_tasks(pid)
    counts = db.signup_counts(pid)
    deliberations = db.list_deliberations(pid)

    lines: List[str] = []
    lines.append(f"# {p['name']} — Groupwork Ultra Agent (Process Log)")
    lines.append("")
    lines.append(f"- Project ID: `{pid}`")
    lines.append(f"- Created at: {p['created_at']}")
    if p["leader"]:
        lines.append(f"- Leader: **{p['leader']}**")
    lines.append(f"- Status: `{p['status']}`")
    lines.append("")

    lines.append("## Ingested documents")
    lines.append("")
    lines.append("### Background")
    if docs_bg:
        for d in docs_bg:
            lines.append(f"- {d['filename']} (sha256={d['sha256'][:10]}…)")  # short
    else:
        lines.append("- (none)")
    lines.append("")
    lines.append("### Assignment requirement")
    if docs_req:
        for d in docs_req:
            lines.append(f"- {d['filename']} (sha256={d['sha256'][:10]}…)")
    else:
        lines.append("- (none)")
    lines.append("")

    lines.append("## Members")
    lines.append("")
    if members:
        for m in members:
            lines.append(f"- {m['name']}")
    else:
        lines.append("- (none)")
    lines.append("")

    if topics:
        lines.append("## Topics")
        lines.append("")
        for t in topics:
            lines.append(f"### {t['title']}")
            lines.append(f"- One-liner: {t['one_liner']}")
            lines.append(f"- Rationale: {t['rationale']}")
            lines.append(f"- Feasibility: {t['feasibility']}")
            ds = json.loads(t["data_sources_json"])
            kq = json.loads(t["key_questions_json"])
            risks = json.loads(t["risks_json"])
            scores = json.loads(t["score_json"] or "{}")
            if ds:
                lines.append("- Data sources: " + "; ".join(ds))
            if kq:
                lines.append("- Key questions: " + "; ".join(kq))
            if risks:
                lines.append("- Risks: " + "; ".join(risks))
            if scores:
                lines.append("- Scores: " + ", ".join([f"{k}={v}" for k, v in scores.items()]))
            lines.append("")

        lines.append("## Voting tally")
        lines.append("")
        if tally:
            for r in tally:
                lines.append(f"- {r['title']}: **{r['votes']}**")
        else:
            lines.append("- (no votes)")
        lines.append("")
        if selected:
            lines.append(f"**Selected topic (current): {selected['title']}**")
            lines.append("")

    if tasks:
        lines.append("## Tasks & signup status")
        lines.append("")
        for task in tasks:
            tid = int(task["id"])
            c = counts.get(tid, 0)
            lines.append(f"### {task['title']} ({c}/{task['max_people']})")
            if task["due_date"]:
                lines.append(f"- Due: {task['due_date']}")
            lines.append(f"- Description: {task['description']}")
            su = db.list_signups_for_task(pid, tid)
            if su:
                lines.append("- Signups: " + ", ".join([x["member_name"] for x in su]))
            lines.append("")

    lines.append("## Deliberation log (compact)")
    lines.append("")
    if deliberations:
        for d in deliberations:
            lines.append(f"- [{d['kind']}] round={d['round_index']} actor={d['actor']} model={d['model']}")
    else:
        lines.append("- (none)")
    lines.append("")

    return "\n".join(lines)


# ----------------------------- Flask app ---------------------------------------

BASE_HTML = """
<!doctype html>
<html lang="en">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{{ title }}</title>
<style>
  body { font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial; margin: 24px; line-height: 1.5; }
  .card { border: 1px solid #ddd; border-radius: 14px; padding: 16px; margin: 12px 0; }
  .muted { color: #666; font-size: 0.95em; }
  .btn { display: inline-block; padding: 10px 14px; border-radius: 10px; border: 1px solid #ccc; text-decoration: none; color: #111; background: #f7f7f7; }
  .btn:hover { background: #f0f0f0; }
  .danger { color: #b00020; }
  label { display:block; margin: 10px 0 6px; font-weight: 600; }
  select { width: 100%; padding: 10px; border-radius: 10px; border: 1px solid #ccc; }
  input[type="radio"] { margin-right: 6px; }
  .radio { margin: 10px 0; }
  .grid { display:grid; grid-template-columns: 1fr; gap: 12px; }
  @media (min-width: 900px) { .grid { grid-template-columns: 1fr 1fr; } }
  pre { white-space: pre-wrap; }
</style>
</head>
<body>
  <h1>{{ title }}</h1>
  <p class="muted">{{ subtitle }}</p>
  {% block content %}{% endblock %}
</body>
</html>
"""


def create_app(db_path: str, pid: str) -> Any:
    if Flask is None:
        raise RuntimeError("Missing dependency: flask. Install with: pip install flask")

    app = Flask(__name__)
    db = DB(db_path)

    @app.get("/")
    def home():
        p = db.get_project(pid)
        members = db.list_members(pid)
        topics = db.list_topics(pid)
        tasks = db.list_tasks(pid)
        counts = db.signup_counts(pid)

        links = []
        for m in members:
            links.append(
                {
                    "name": m["name"],
                    "vote": url_for("vote", token=m["token"], _external=True),
                    "signup": url_for("signup", token=m["token"], _external=True),
                }
            )

        final_md = db.get_artifact(pid, "final_markdown")

        html = """
        {% extends base %}
        {% block content %}
          <div class="card">
            <div><b>Project ID:</b> {{ pid }}</div>
            <div><b>Leader:</b> {{ leader or "(not set)" }}</div>
            <div><b>Status:</b> {{ status }}</div>
          </div>

          <div class="grid">
            <div class="card">
              <h2>Topics</h2>
              {% if topics|length == 0 %}
                <p class="muted">No topics yet. Generate them via CLI.</p>
              {% else %}
                {% for t in topics %}
                  <div style="margin: 12px 0;">
                    <b>#{{ loop.index }} {{ t.title }}</b>
                    <div class="muted">{{ t.one_liner }}</div>
                  </div>
                {% endfor %}
              {% endif %}
              <p><a class="btn" href="{{ url_for('tally') }}">View vote tally</a></p>
            </div>

            <div class="card">
              <h2>Tasks</h2>
              {% if tasks|length == 0 %}
                <p class="muted">No tasks yet. Generate them via CLI.</p>
              {% else %}
                {% for task in tasks %}
                  <div style="margin: 12px 0;">
                    <b>{{ task.title }}</b>
                    <div class="muted">{{ counts.get(task.id, 0) }}/{{ task.max_people }}</div>
                  </div>
                {% endfor %}
              {% endif %}
              <p><a class="btn" href="{{ url_for('tasks') }}">View tasks</a></p>
            </div>
          </div>

          <div class="card">
            <h2>Member Links</h2>
            <p class="muted">Send each person their unique link to vote / sign up.</p>
            {% if links|length == 0 %}
              <p class="muted">(No members set)</p>
            {% else %}
              <ul>
                {% for x in links %}
                  <li><b>{{ x.name }}</b> — <a href="{{ x.vote }}">Vote</a> | <a href="{{ x.signup }}">Signup</a></li>
                {% endfor %}
              </ul>
            {% endif %}
          </div>

          <div class="card">
            <h2>Exports</h2>
            <p><a class="btn" href="{{ url_for('export_process') }}">Export process Markdown</a></p>
            {% if final_md %}
              <p><a class="btn" href="{{ url_for('export_final') }}">Export final Markdown</a></p>
            {% else %}
              <p class="muted">Final Markdown not generated yet.</p>
            {% endif %}
          </div>
        {% endblock %}
        """
        return render_template_string(
            html,
            base=BASE_HTML,
            title=p["name"],
            subtitle="Local Web (vote & signup) — Ultra Agent",
            pid=pid,
            leader=p["leader"],
            status=p["status"],
            topics=topics,
            tasks=tasks,
            counts=counts,
            links=links,
            final_md=final_md,
        )

    @app.get("/tally")
    def tally():
        p = db.get_project(pid)
        rows = db.tally_votes(pid)
        selected = db.selected_topic(pid)
        html = """
        {% extends base %}
        {% block content %}
          <div class="card">
            <h2>Vote tally</h2>
            {% if rows|length == 0 %}
              <p class="muted">No topics or votes yet.</p>
            {% else %}
              <ul>
                {% for r in rows %}
                  <li>{{ r.title }} — <b>{{ r.votes }}</b></li>
                {% endfor %}
              </ul>
              {% if selected %}
                <p><b>Current selected:</b> {{ selected.title }}</p>
              {% endif %}
            {% endif %}
            <p><a class="btn" href="{{ url_for('home') }}">Back</a></p>
          </div>
        {% endblock %}
        """
        return render_template_string(html, base=BASE_HTML, title=p["name"], subtitle="Vote tally", rows=rows, selected=selected)

    @app.get("/tasks")
    def tasks():
        p = db.get_project(pid)
        tasks = db.list_tasks(pid)
        counts = db.signup_counts(pid)
        details = []
        for t in tasks:
            su = db.list_signups_for_task(pid, int(t["id"]))
            details.append((t, su))
        html = """
        {% extends base %}
        {% block content %}
          <div class="card">
            <h2>Tasks</h2>
            {% if tasks|length == 0 %}
              <p class="muted">No tasks.</p>
            {% else %}
              {% for t, su in details %}
                <div class="card">
                  <div><b>{{ t.title }}</b> <span class="muted">({{ counts.get(t.id, 0) }}/{{ t.max_people }})</span></div>
                  {% if t.due_date %}<div class="muted">Due: {{ t.due_date }}</div>{% endif %}
                  <div style="margin-top:8px;">{{ t.description }}</div>
                  {% if su|length > 0 %}
                    <div class="muted" style="margin-top:8px;">Signed: {{ su | map(attribute='member_name') | join(', ') }}</div>
                  {% endif %}
                </div>
              {% endfor %}
            {% endif %}
            <p><a class="btn" href="{{ url_for('home') }}">Back</a></p>
          </div>
        {% endblock %}
        """
        return render_template_string(html, base=BASE_HTML, title=p["name"], subtitle="Tasks", tasks=tasks, counts=counts, details=details)

    @app.route("/vote/<token>", methods=["GET", "POST"])
    def vote(token: str):
        p = db.get_project(pid)
        member = db.member_by_token(pid, token)
        if not member:
            abort(404)
        topics = db.list_topics(pid)
        if not topics:
            return render_template_string(
                """{% extends base %}{% block content %}
                <div class="card"><p class="danger">No topics available.</p>
                <p><a class="btn" href="{{ url_for('home') }}">Back</a></p></div>
                {% endblock %}""",
                base=BASE_HTML,
                title=p["name"],
                subtitle="Vote",
            )

        # randomize display order per hour per token
        ids = [int(t["id"]) for t in topics]
        rng = random.Random(token + ":" + dt.datetime.utcnow().strftime("%Y%m%d%H"))
        rng.shuffle(ids)
        topic_map = {int(t["id"]): t for t in topics}
        ordered = [topic_map[i] for i in ids]

        if request.method == "POST":
            choice = request.form.get("topic_id", "").strip()
            if not choice.isdigit():
                abort(400)
            topic_id = int(choice)
            if not db.topic_by_id(pid, topic_id):
                abort(400)
            db.cast_vote(pid, int(member["id"]), topic_id)
            return redirect(url_for("tally"))

        html = """
        {% extends base %}
        {% block content %}
          <div class="card">
            <p>Hi <b>{{ member.name }}</b>. Pick one topic. You can re-submit to overwrite.</p>
            <form method="post">
              {% for t in topics %}
                <div class="radio">
                  <label>
                    <input type="radio" name="topic_id" value="{{ t.id }}" required>
                    <b>{{ t.title }}</b> <span class="muted">— {{ t.one_liner }}</span>
                  </label>
                </div>
              {% endfor %}
              <button class="btn" type="submit">Submit vote</button>
              <a class="btn" href="{{ url_for('home') }}">Back</a>
            </form>
          </div>
        {% endblock %}
        """
        return render_template_string(html, base=BASE_HTML, title=p["name"], subtitle="Vote", member=member, topics=ordered)

    @app.route("/signup/<token>", methods=["GET", "POST"])
    def signup(token: str):
        p = db.get_project(pid)
        member = db.member_by_token(pid, token)
        if not member:
            abort(404)
        tasks = db.list_tasks(pid)
        if not tasks:
            return render_template_string(
                """{% extends base %}{% block content %}
                <div class="card"><p class="danger">No tasks available.</p>
                <p><a class="btn" href="{{ url_for('home') }}">Back</a></p></div>
                {% endblock %}""",
                base=BASE_HTML,
                title=p["name"],
                subtitle="Signup",
            )
        counts = db.signup_counts(pid)
        current = db.member_current_signup(pid, int(member["id"]))

        if request.method == "POST":
            task_id = request.form.get("task_id", "").strip()
            if not task_id.isdigit():
                abort(400)
            ok, msg = db.signup_task(pid, int(member["id"]), int(task_id), allow_change=True)
            if ok:
                return redirect(url_for("tasks"))
            return render_template_string(
                """{% extends base %}{% block content %}
                <div class="card"><p class="danger">{{ msg }}</p>
                <p><a class="btn" href="{{ url_for('signup', token=token) }}">Try again</a></p></div>
                {% endblock %}""",
                base=BASE_HTML,
                title=p["name"],
                subtitle="Signup failed",
                msg=msg,
                token=token,
            )

        html = """
        {% extends base %}
        {% block content %}
          <div class="card">
            <p>Hi <b>{{ member.name }}</b>. Choose one task (can change once; full tasks disabled).</p>
            {% if current %}
              <p class="muted">Current signed task ID: {{ current }}</p>
            {% endif %}
            <form method="post">
              <label>Select task</label>
              <select name="task_id" required>
                {% for t in tasks %}
                  {% set c = counts.get(t.id, 0) %}
                  {% set full = (c >= t.max_people) %}
                  <option value="{{ t.id }}" {% if full %}disabled{% endif %}>
                    [{{ c }}/{{ t.max_people }}] {{ t.title }} {% if full %}(FULL){% endif %}
                  </option>
                {% endfor %}
              </select>
              <div style="margin-top:12px;">
                <button class="btn" type="submit">Confirm</button>
                <a class="btn" href="{{ url_for('home') }}">Back</a>
              </div>
            </form>
          </div>
        {% endblock %}
        """
        return render_template_string(html, base=BASE_HTML, title=p["name"], subtitle="Signup", member=member, tasks=tasks, counts=counts, current=current)

    @app.get("/export_process.md")
    def export_process():
        md = export_process_markdown(db, pid)
        return app.response_class(md, mimetype="text/markdown; charset=utf-8")

    @app.get("/export_final.md")
    def export_final():
        md = db.get_artifact(pid, "final_markdown") or "(final markdown not generated)"
        return app.response_class(md, mimetype="text/markdown; charset=utf-8")

    return app


# ----------------------------- CLI Commands ------------------------------------

def parse_members(members_csv: str) -> List[str]:
    parts = [x.strip() for x in members_csv.replace("，", ",").split(",")]
    return [p for p in parts if p]


def cmd_init(args: argparse.Namespace) -> None:
    db = DB(args.db)

    settings = {
        "api_base": args.api_base,
        "models": {
            "leader": args.leader_model,
            "critic": args.critic_model,
            "researcher": args.researcher_model,
            "methodologist": args.methodologist_model,
            "editor": args.editor_model,
            "redteam": args.redteam_model,
            "controller": args.controller_model,
        },
        "temperature": args.temperature,
        "max_tokens": args.max_tokens,
        "openrouter_headers": {
            "HTTP-Referer": args.http_referer or "",
            "X-Title": args.app_title or "",
        },
        "bg_dir": args.background_dir,
        "req_dir": args.assignment_dir,
    }

    pid = db.create_project(args.name, settings=settings)
    added = db.add_members(pid, parse_members(args.members) if args.members else [])

    print(f"✅ Project created: {pid}\n")
    if added:
        print("Members:")
        for name, token in added:
            print(f"- {name}: {token[:8]}…")
    else:
        print("Members: (none)")
    print("\nNext recommended steps:")
    print(f"  1) ingest docs:     python {os.path.basename(__file__)} ingest {pid}")
    print(f"  2) propose topics:  python {os.path.basename(__file__)} propose-topics {pid} --mode autopilot")
    print(f"  3) tasks:           python {os.path.basename(__file__)} generate-tasks {pid} --mode autopilot")
    print(f"  4) final markdown:  python {os.path.basename(__file__)} generate-final {pid}")
    print(f"  5) serve web:       python {os.path.basename(__file__)} serve {pid} --host 127.0.0.1 --port 8080")


def cmd_ingest(args: argparse.Namespace) -> None:
    db = DB(args.db)
    p = db.get_project(args.project_id)
    settings = json.loads(p["settings_json"] or "{}")
    bg_dir = args.background_dir or settings.get("bg_dir") or DEFAULT_BG_DIR
    req_dir = args.assignment_dir or settings.get("req_dir") or DEFAULT_REQ_DIR

    # Ingest background
    bg_docs = load_directory_documents(bg_dir)
    req_docs = load_directory_documents(req_dir)

    if args.clear:
        db.clear_documents(args.project_id)

    def ingest_kind(kind: str, docs: List[Tuple[str, str]]) -> None:
        for filename, text in docs:
            doc_id = db.insert_document(args.project_id, kind, filename, text)
            chunks = chunk_text(text)
            db.insert_chunks(args.project_id, doc_id, chunks)

    ingest_kind("background", bg_docs)
    ingest_kind("requirement", req_docs)

    db.update_project(args.project_id, status="ingested")
    print(f"✅ Ingested: background={len(bg_docs)} files, requirement={len(req_docs)} files.")


def _build_llm_from_project(db: DB, pid: str, args: argparse.Namespace) -> Tuple[LLMClient, Dict[str, str], Dict[str, Any]]:
    p = db.get_project(pid)
    settings = json.loads(p["settings_json"] or "{}")

    api_base = args.api_base or settings.get("api_base") or DEFAULT_API_BASE
    temperature = args.temperature if args.temperature is not None else settings.get("temperature", 0.4)
    max_tokens = args.max_tokens if args.max_tokens is not None else settings.get("max_tokens", 2800)

    models = settings.get("models", {}) or {}
    merged_models = {
        "leader": args.leader_model or models.get("leader") or DEFAULT_MODELS["leader"],
        "critic": args.critic_model or models.get("critic") or DEFAULT_MODELS["critic"],
        "researcher": args.researcher_model or models.get("researcher") or DEFAULT_MODELS["researcher"],
        "methodologist": args.methodologist_model or models.get("methodologist") or DEFAULT_MODELS["methodologist"],
        "editor": args.editor_model or models.get("editor") or DEFAULT_MODELS["editor"],
        "redteam": args.redteam_model or models.get("redteam") or DEFAULT_MODELS["redteam"],
        "controller": args.controller_model or models.get("controller") or DEFAULT_MODELS["controller"],
    }

    # Extra headers (OpenRouter friendly)
    openrouter_headers = settings.get("openrouter_headers", {}) or {}
    extra_headers = {}
    if args.http_referer is not None:
        extra_headers["HTTP-Referer"] = args.http_referer
    elif openrouter_headers.get("HTTP-Referer"):
        extra_headers["HTTP-Referer"] = openrouter_headers["HTTP-Referer"]
    if args.app_title is not None:
        extra_headers["X-Title"] = args.app_title
    elif openrouter_headers.get("X-Title"):
        extra_headers["X-Title"] = openrouter_headers["X-Title"]

    llm = LLMClient(
        LLMCallConfig(
            api_base=api_base,
            api_key=choose_api_key(args.api_key),
            temperature=float(temperature),
            max_tokens=int(max_tokens),
            extra_headers=extra_headers,
        )
    )
    return llm, merged_models, settings


def cmd_pick_leader(args: argparse.Namespace) -> None:
    db = DB(args.db)
    members = [m["name"] for m in db.list_members(args.project_id)]
    if not members:
        raise SystemExit("No members in this project. Re-run init with --members ...")

    if args.method in ("lottery", "dice"):
        leader, meta = pick_leader_commit_reveal(members, method=args.method, seed=args.seed)
        db.update_project(args.project_id, leader=leader)
        print(f"🎲 Leader: {leader}")
        print("Fairness note (commit-reveal):")
        print(f"- seed SHA256 commit: {meta['commit_sha256']}")
        print(f"- seed (reveal): {meta['seed']}")
        if "rolls" in meta:
            print("- dice rolls:")
            for k, v in meta["rolls"].items():
                print(f"  - {k}: {v}")
        return

    # AI method
    llm, models, _ = _build_llm_from_project(db, args.project_id, args)
    decision = pick_leader_ai(llm, models["controller"], members)
    leader = decision.get("leader")
    if leader not in members:
        leader = members[0]
    db.update_project(args.project_id, leader=leader)
    print(f"🤖 AI-picked leader: {leader}")
    print(f"- reason: {decision.get('reason','(none)')}")


def cmd_propose_topics(args: argparse.Namespace) -> None:
    db = DB(args.db)
    _ = db.get_project(args.project_id)

    llm, models, _settings = _build_llm_from_project(db, args.project_id, args)
    agents = default_agents(models)

    # Ensure ingestion exists
    if not db.list_documents(args.project_id):
        print("[warn] No ingested documents found. Run `ingest` first for best results.")

    topics = deliberate_topics(
        db=db,
        llm=llm,
        pid=args.project_id,
        agents=agents,
        controller_model=models["controller"],
        n_topics=args.n,
        min_rounds=args.min_rounds,
        max_rounds=args.max_rounds,
        rag_k=args.rag_k,
    )

    db.clear_topics(args.project_id)
    db.insert_topics(args.project_id, topics)
    db.update_project(args.project_id, status="topics_ready")

    print(f"✅ Inserted {len(topics)} topics into DB.")

    if args.mode == "autopilot":
        pick = auto_pick_best_topic(llm, models["controller"], topics)
        sel_title = str(pick.get("selected_title", "")).strip()
        # mark selection via "leader" field or status; selection is otherwise vote-based
        db.update_project(args.project_id, status=f"topics_autopicked:{sel_title}")
        print(f"🤖 Autopilot selected: {sel_title}")
        print(f"- reason: {pick.get('reason','(none)')}")
    else:
        print("Human mode: start server and send Vote links.")
        print(f"  python {os.path.basename(__file__)} serve {args.project_id} --host 0.0.0.0 --port 8000")


def cmd_generate_tasks(args: argparse.Namespace) -> None:
    db = DB(args.db)
    p = db.get_project(args.project_id)
    llm, models, _settings = _build_llm_from_project(db, args.project_id, args)
    agents = default_agents(models)

    topics_rows = db.list_topics(args.project_id)
    if not topics_rows:
        raise SystemExit("No topics found. Run propose-topics first.")

    selected = None
    if args.mode == "human":
        selected = db.selected_topic(args.project_id)
        if not selected:
            raise SystemExit("No selected topic via votes yet. Run server, vote, then retry.")
        selected_topic = {
            "title": selected["title"],
            "one_liner": selected["one_liner"],
            "rationale": selected["rationale"],
            "feasibility": selected["feasibility"],
            "data_sources": json.loads(selected["data_sources_json"]),
            "key_questions": json.loads(selected["key_questions_json"]),
            "risks": json.loads(selected["risks_json"]),
            "scores": json.loads(selected["score_json"] or "{}"),
        }
    else:
        # autopilot picks best
        topics = []
        for r in topics_rows:
            topics.append({
                "title": r["title"],
                "one_liner": r["one_liner"],
                "rationale": r["rationale"],
                "feasibility": r["feasibility"],
                "data_sources": json.loads(r["data_sources_json"]),
                "key_questions": json.loads(r["key_questions_json"]),
                "risks": json.loads(r["risks_json"]),
                "scores": json.loads(r["score_json"] or "{}"),
            })
        pick = auto_pick_best_topic(llm, models["controller"], topics)
        sel_title = str(pick.get("selected_title", "")).strip()
        # find it
        selected_topic = next((t for t in topics if t["title"] == sel_title), topics[0])
        db.update_project(args.project_id, status=f"topic_autopicked:{selected_topic['title']}")
        print(f"🤖 Autopilot selected topic: {selected_topic['title']}")

    member_count = len(db.list_members(args.project_id)) or 3  # assume "virtual group size" if no members

    tasks = deliberate_tasks(
        db=db,
        llm=llm,
        pid=args.project_id,
        agents=agents,
        controller_model=models["controller"],
        topic_title=selected_topic["title"],
        member_count=member_count,
        min_rounds=args.min_rounds,
        max_rounds=args.max_rounds,
        rag_k=args.rag_k,
    )

    db.clear_tasks(args.project_id)
    db.insert_tasks(args.project_id, tasks)
    db.update_project(args.project_id, status="tasks_ready")

    print(f"✅ Inserted {len(tasks)} tasks.")
    if args.mode == "human":
        print("Human mode: send each member their Signup link (home page shows them).")
    else:
        print("Autopilot: tasks created (you can still use web signup optionally).")


def cmd_generate_final(args: argparse.Namespace) -> None:
    db = DB(args.db)
    _ = db.get_project(args.project_id)
    llm, models, _settings = _build_llm_from_project(db, args.project_id, args)
    agents = default_agents(models)

    topics_rows = db.list_topics(args.project_id)
    if not topics_rows:
        raise SystemExit("No topics. Run propose-topics first.")
    tasks_rows = db.list_tasks(args.project_id)
    if not tasks_rows:
        raise SystemExit("No tasks. Run generate-tasks first.")

    # Decide selected topic
    selected = db.selected_topic(args.project_id)
    if selected:
        selected_topic = {
            "title": selected["title"],
            "one_liner": selected["one_liner"],
            "rationale": selected["rationale"],
            "feasibility": selected["feasibility"],
            "data_sources": json.loads(selected["data_sources_json"]),
            "key_questions": json.loads(selected["key_questions_json"]),
            "risks": json.loads(selected["risks_json"]),
            "scores": json.loads(selected["score_json"] or "{}"),
        }
    else:
        # autopick if no votes
        topics = []
        for r in topics_rows:
            topics.append({
                "title": r["title"],
                "one_liner": r["one_liner"],
                "rationale": r["rationale"],
                "feasibility": r["feasibility"],
                "data_sources": json.loads(r["data_sources_json"]),
                "key_questions": json.loads(r["key_questions_json"]),
                "risks": json.loads(r["risks_json"]),
                "scores": json.loads(r["score_json"] or "{}"),
            })
        pick = auto_pick_best_topic(llm, models["controller"], topics)
        sel_title = str(pick.get("selected_title", "")).strip()
        selected_topic = next((t for t in topics if t["title"] == sel_title), topics[0])

    tasks = []
    for t in tasks_rows:
        tasks.append({"title": t["title"], "description": t["description"], "max_people": int(t["max_people"]), "due_date": t["due_date"]})

    md = deliberate_final_markdown(
        db=db,
        llm=llm,
        pid=args.project_id,
        agents=agents,
        controller_model=models["controller"],
        selected_topic=selected_topic,
        tasks=tasks,
        min_rounds=args.min_rounds,
        max_rounds=args.max_rounds,
        rag_k=args.rag_k,
    )
    db.update_project(args.project_id, status="final_ready")

    out_path = args.out or "Final_Group_Report.md"
    with open(out_path, "w", encoding="utf-8") as f:
        f.write(md)
    print(f"✅ Final Markdown written to: {out_path}")


def cmd_export(args: argparse.Namespace) -> None:
    db = DB(args.db)
    _ = db.get_project(args.project_id)

    if args.kind == "process":
        md = export_process_markdown(db, args.project_id)
        if args.out:
            with open(args.out, "w", encoding="utf-8") as f:
                f.write(md)
            print(f"✅ Exported process markdown to {args.out}")
        else:
            print(md)
        return

    if args.kind == "final":
        md = db.get_artifact(args.project_id, "final_markdown") or "(no final markdown yet)"
        if args.out:
            with open(args.out, "w", encoding="utf-8") as f:
                f.write(md)
            print(f"✅ Exported final markdown to {args.out}")
        else:
            print(md)
        return

    raise SystemExit("Unknown export kind.")


def cmd_serve(args: argparse.Namespace) -> None:
    if Flask is None:
        raise SystemExit("Missing dependency: flask. Install with: pip install flask")
    app = create_app(args.db, args.project_id)
    app.run(host=args.host, port=args.port, debug=args.debug)


def cmd_run_all(args: argparse.Namespace) -> None:
    """
    End-to-end pipeline (Autopilot):
      ingest -> propose-topics(autopilot) -> generate-tasks(autopilot) -> generate-final
    """
    # Ingest
    cmd_ingest(argparse.Namespace(
        db=args.db,
        project_id=args.project_id,
        clear=args.clear_ingest,
        background_dir=args.background_dir,
        assignment_dir=args.assignment_dir,
    ))
    # topics
    cmd_propose_topics(argparse.Namespace(
        db=args.db,
        project_id=args.project_id,
        api_base=args.api_base,
        api_key=args.api_key,
        leader_model=args.leader_model,
        critic_model=args.critic_model,
        researcher_model=args.researcher_model,
        methodologist_model=args.methodologist_model,
        editor_model=args.editor_model,
        redteam_model=args.redteam_model,
        controller_model=args.controller_model,
        temperature=args.temperature,
        max_tokens=args.max_tokens,
        http_referer=args.http_referer,
        app_title=args.app_title,
        n=args.n,
        min_rounds=args.min_rounds,
        max_rounds=args.max_rounds,
        rag_k=args.rag_k,
        mode="autopilot",
    ))
    # tasks
    cmd_generate_tasks(argparse.Namespace(
        db=args.db,
        project_id=args.project_id,
        api_base=args.api_base,
        api_key=args.api_key,
        leader_model=args.leader_model,
        critic_model=args.critic_model,
        researcher_model=args.researcher_model,
        methodologist_model=args.methodologist_model,
        editor_model=args.editor_model,
        redteam_model=args.redteam_model,
        controller_model=args.controller_model,
        temperature=args.temperature,
        max_tokens=args.max_tokens,
        http_referer=args.http_referer,
        app_title=args.app_title,
        min_rounds=args.min_rounds,
        max_rounds=args.max_rounds,
        rag_k=args.rag_k,
        mode="autopilot",
    ))
    # final
    cmd_generate_final(argparse.Namespace(
        db=args.db,
        project_id=args.project_id,
        api_base=args.api_base,
        api_key=args.api_key,
        leader_model=args.leader_model,
        critic_model=args.critic_model,
        researcher_model=args.researcher_model,
        methodologist_model=args.methodologist_model,
        editor_model=args.editor_model,
        redteam_model=args.redteam_model,
        controller_model=args.controller_model,
        temperature=args.temperature,
        max_tokens=args.max_tokens,
        http_referer=args.http_referer,
        app_title=args.app_title,
        min_rounds=args.min_rounds,
        max_rounds=args.max_rounds,
        rag_k=args.rag_k,
        out=args.out,
    ))


# ----------------------------- Parser ------------------------------------------

def build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(
        prog="groupwork_ultra_agent",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
        description="Groupwork Ultra Agent: multi-model deliberation + local web voting/signup + RAG + doc ingestion.",
    )
    p.add_argument("--db", default=DEFAULT_DB, help="SQLite db path")

    # global LLM config overrides
    p.add_argument("--api-base", default=None, help=f"OpenAI-compatible API base (default project setting or {DEFAULT_API_BASE})")
    p.add_argument("--api-key", default=None, help="API key (or env: LLM_API_KEY / OPENROUTER_API_KEY / OPENAI_API_KEY)")
    p.add_argument("--temperature", type=float, default=None, help="Sampling temperature")
    p.add_argument("--max-tokens", type=int, default=None, help="Max tokens per completion")
    p.add_argument("--http-referer", default=None, help="Optional OpenRouter HTTP-Referer header")
    p.add_argument("--app-title", default=None, help="Optional OpenRouter X-Title header")

    # model overrides
    p.add_argument("--leader-model", default=None)
    p.add_argument("--critic-model", default=None)
    p.add_argument("--researcher-model", default=None)
    p.add_argument("--methodologist-model", default=None)
    p.add_argument("--editor-model", default=None)
    p.add_argument("--redteam-model", default=None)
    p.add_argument("--controller-model", default=None)

    sub = p.add_subparsers(dest="cmd", required=True)

    sp = sub.add_parser("init", help="Create a new project")
    sp.add_argument("--name", required=True, help="Project / course name")
    sp.add_argument("--members", default="", help="Comma-separated member names (optional)")
    sp.add_argument("--background-dir", default=DEFAULT_BG_DIR, help="Folder for background docs")
    sp.add_argument("--assignment-dir", default=DEFAULT_REQ_DIR, help="Folder for requirement docs")
    sp.add_argument("--api-base", default=DEFAULT_API_BASE)
    sp.add_argument("--leader-model", default=DEFAULT_MODELS["leader"])
    sp.add_argument("--critic-model", default=DEFAULT_MODELS["critic"])
    sp.add_argument("--researcher-model", default=DEFAULT_MODELS["researcher"])
    sp.add_argument("--methodologist-model", default=DEFAULT_MODELS["methodologist"])
    sp.add_argument("--editor-model", default=DEFAULT_MODELS["editor"])
    sp.add_argument("--redteam-model", default=DEFAULT_MODELS["redteam"])
    sp.add_argument("--controller-model", default=DEFAULT_MODELS["controller"])
    sp.add_argument("--temperature", type=float, default=0.4)
    sp.add_argument("--max-tokens", type=int, default=2800)
    sp.add_argument("--http-referer", default=None)
    sp.add_argument("--app-title", default=None)
    sp.set_defaults(func=cmd_init)

    sp = sub.add_parser("ingest", help="Ingest documents from folders into DB + chunks + FTS")
    sp.add_argument("project_id")
    sp.add_argument("--clear", action="store_true", help="Clear old ingested documents first")
    sp.add_argument("--background-dir", default=None)
    sp.add_argument("--assignment-dir", default=None)
    sp.set_defaults(func=cmd_ingest)

    sp = sub.add_parser("pick-leader", help="Pick leader: dice/lottery (commit-reveal) or ai (smart pick)")
    sp.add_argument("project_id")
    sp.add_argument("--method", choices=["dice", "lottery", "ai"], default="dice")
    sp.add_argument("--seed", default=None)
    sp.set_defaults(func=cmd_pick_leader)

    sp = sub.add_parser("propose-topics", help="Deliberate topics with multi-agent loops; store in DB")
    sp.add_argument("project_id")
    sp.add_argument("--n", type=int, default=8, help="Number of topics (3..12)")
    sp.add_argument("--min-rounds", type=int, default=2)
    sp.add_argument("--max-rounds", type=int, default=6)
    sp.add_argument("--rag-k", type=int, default=8, help="Top-k chunks for RAG")
    sp.add_argument("--mode", choices=["autopilot", "human"], default="autopilot")
    sp.set_defaults(func=cmd_propose_topics)

    sp = sub.add_parser("generate-tasks", help="Deliberate tasks based on selected topic; store in DB")
    sp.add_argument("project_id")
    sp.add_argument("--min-rounds", type=int, default=2)
    sp.add_argument("--max-rounds", type=int, default=6)
    sp.add_argument("--rag-k", type=int, default=8)
    sp.add_argument("--mode", choices=["autopilot", "human"], default="autopilot")
    sp.set_defaults(func=cmd_generate_tasks)

    sp = sub.add_parser("generate-final", help="Deliberate the final comprehensive Markdown and save artifact")
    sp.add_argument("project_id")
    sp.add_argument("--min-rounds", type=int, default=2)
    sp.add_argument("--max-rounds", type=int, default=6)
    sp.add_argument("--rag-k", type=int, default=10)
    sp.add_argument("--out", default="Final_Group_Report.md")
    sp.set_defaults(func=cmd_generate_final)

    sp = sub.add_parser("export", help="Export markdown (process or final)")
    sp.add_argument("project_id")
    sp.add_argument("--kind", choices=["process", "final"], default="process")
    sp.add_argument("--out", default=None)
    sp.set_defaults(func=cmd_export)

    sp = sub.add_parser("serve", help="Run local web server (vote/signup/exports)")
    sp.add_argument("project_id")
    sp.add_argument("--host", default=DEFAULT_HOST)
    sp.add_argument("--port", type=int, default=DEFAULT_PORT)
    sp.add_argument("--debug", action="store_true")
    sp.set_defaults(func=cmd_serve)

    sp = sub.add_parser("run-all", help="Autopilot: ingest -> topics -> tasks -> final")
    sp.add_argument("project_id")
    sp.add_argument("--clear-ingest", action="store_true")
    sp.add_argument("--background-dir", default=None)
    sp.add_argument("--assignment-dir", default=None)
    sp.add_argument("--n", type=int, default=8)
    sp.add_argument("--min-rounds", type=int, default=2)
    sp.add_argument("--max-rounds", type=int, default=6)
    sp.add_argument("--rag-k", type=int, default=8)
    sp.add_argument("--out", default="Final_Group_Report.md")
    sp.set_defaults(func=cmd_run_all)

    return p


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()

    # fill model defaults if missing
    args.leader_model = args.leader_model or DEFAULT_MODELS["leader"]
    args.critic_model = args.critic_model or DEFAULT_MODELS["critic"]
    args.researcher_model = args.researcher_model or DEFAULT_MODELS["researcher"]
    args.methodologist_model = args.methodologist_model or DEFAULT_MODELS["methodologist"]
    args.editor_model = args.editor_model or DEFAULT_MODELS["editor"]
    args.redteam_model = args.redteam_model or DEFAULT_MODELS["redteam"]
    args.controller_model = args.controller_model or DEFAULT_MODELS["controller"]

    args.func(args)


if __name__ == "__main__":
    main()
